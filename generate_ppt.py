"""
LifeOS Healthcare AI — PowerPoint Presentation Generator (Light Theme)
======================================================================
Run:
    pip install python-pptx
    python generate_ppt.py

Output: LifeOS_Healthcare_AI_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────── BRAND COLORS ────────────────────────────
PRIMARY       = RGBColor(0x25, 0x63, 0xEB)   # #2563EB
PRIMARY_DARK  = RGBColor(0x1D, 0x4E, 0xD8)   # #1D4ED8
PRIMARY_LIGHT = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6
SECONDARY     = RGBColor(0x06, 0xB6, 0xD4)   # #06B6D4 (Cyan)
ACCENT        = RGBColor(0x25, 0x63, 0xEB)   # Use primary blue as accent on light bg
SUCCESS       = RGBColor(0x16, 0xA3, 0x4A)   # #16A34A  (slightly darker green for contrast)
WARNING       = RGBColor(0xD9, 0x77, 0x06)   # #D97706  (darker amber for contrast)
DANGER        = RGBColor(0xDC, 0x26, 0x26)   # #DC2626
PURPLE        = RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED

# ─── Light Theme Backgrounds (from DESIGN.md light mode) ───
SLIDE_BG      = RGBColor(0xF0, 0xF4, 0xFF)   # #F0F4FF  --bg-color
PANEL_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF  --panel-bg
CARD_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
BG_SECONDARY  = RGBColor(0xF1, 0xF5, 0xF9)   # #F1F5F9  --bg-secondary
BG_SUBTLE     = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0  for alternating rows

# ─── Light Theme Text ───
TEXT_MAIN     = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  --text-main
TEXT_BODY     = RGBColor(0x33, 0x41, 0x55)   # #334155  secondary body text
TEXT_MUTED    = RGBColor(0x64, 0x74, 0x8B)   # #64748B  muted (darker than dark-mode muted)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

# ─── Decorative ───
CIRCLE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)   # #DBEAFE  soft blue circle
CIRCLE_CYAN   = RGBColor(0xCC, 0xFB, 0xF1)   # #CCFBF1  soft cyan circle
CARD_BORDER   = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0  subtle border

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ──────────────────────────── HELPERS ────────────────────────────
def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    """Add a thin accent/divider bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_MAIN, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_BODY, icon="▸", font_name="Calibri"):
    """Add a bullet list with custom icon."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             card_color=CARD_BG, title_color=PRIMARY, body_color=TEXT_BODY,
             border_color=CARD_BORDER):
    """Add a card with a title and bullet points."""
    add_shape(slide, left, top, width, height, card_color, border_color=border_color, border_width=1)
    # Accent bar at top of card
    add_accent_bar(slide, left + Inches(0.3), top + Inches(0.15), Inches(0.6), Pt(4), title_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                    height - Inches(0.9), body_lines, font_size=13, color=body_color, icon="•")


def add_section_header_shape(slide, text, subtitle=""):
    """Add a left-aligned big section header."""
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 text, font_size=36, color=TEXT_MAIN, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.2), Pt(5), PRIMARY)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
                     subtitle, font_size=18, color=TEXT_MUTED)


def add_page_number(slide, num, total):
    """Add subtle page number at bottom right."""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=TEXT_MUTED,
                 alignment=PP_ALIGN.RIGHT)


# ──────────────────────────── SLIDE BUILDERS ────────────────────────────

def build_title_slide(prs):
    """Slide 1: Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, SLIDE_BG)

    # Large decorative circle (top-right) — soft blue
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CIRCLE_LIGHT
    circle.line.fill.background()

    # Small decorative circle (bottom-left) — soft cyan
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5.5), Inches(3.5), Inches(3.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = CIRCLE_CYAN
    circle2.line.fill.background()

    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(10), Inches(1.2),
                 "LifeOS", font_size=60, color=TEXT_MAIN, bold=True)
    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.9), Inches(10), Inches(0.8),
                 "AI-Powered Healthcare Operating System", font_size=30,
                 color=PRIMARY, bold=False)
    # Accent line
    add_accent_bar(slide, Inches(1), Inches(3.7), Inches(2), Pt(5), SECONDARY)
    # Description
    add_text_box(slide, Inches(1), Inches(4.0), Inches(8), Inches(1.2),
                 "A production-ready full-stack application with 80+ REST API endpoints,\n"
                 "16 healthcare modules, and deep AI integration — built with FastAPI & React.",
                 font_size=17, color=TEXT_MUTED)
    # Bottom badge
    add_shape(slide, Inches(1), Inches(5.8), Inches(3.8), Inches(0.5), WHITE,
              border_color=PRIMARY, border_width=1)
    add_text_box(slide, Inches(1.15), Inches(5.82), Inches(3.5), Inches(0.45),
                 "⚡  FastAPI  ·  React  ·  PostgreSQL  ·  Groq AI",
                 font_size=13, color=PRIMARY, alignment=PP_ALIGN.CENTER)


def build_overview_slide(prs):
    """Slide 2: Project Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "Project Overview", "What is LifeOS?")

    # Description card
    add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.4), WHITE,
              border_color=CARD_BORDER, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(2.15), Inches(11), Inches(1.2),
                 "LifeOS is a comprehensive AI-powered healthcare platform that serves as a personal health "
                 "operating system. It combines medical record management, health tracking, AI-driven insights, "
                 "and gamification to help users take control of their health journey.",
                 font_size=17, color=TEXT_BODY)

    # Key stats cards
    stats = [
        ("80+",  "API Endpoints",   "Across 16 modules",    PRIMARY),
        ("17",   "API Routers",     "RESTful architecture",  SECONDARY),
        ("11",   "Database Models", "SQLAlchemy ORM",        SUCCESS),
        ("5",    "AI Modules",      "Groq Llama 3.3 70B",   PURPLE),
    ]
    x_start = Inches(0.8)
    card_w = Inches(2.6)
    gap = Inches(0.27)
    for i, (num, label, desc, color) in enumerate(stats):
        x = x_start + i * (card_w + gap)
        add_shape(slide, x, Inches(3.8), card_w, Inches(2.6), WHITE,
                  border_color=color, border_width=2)
        add_text_box(slide, x + Inches(0.3), Inches(4.0), card_w - Inches(0.6), Inches(0.8),
                     num, font_size=44, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(4.9), card_w - Inches(0.6), Inches(0.5),
                     label, font_size=17, color=TEXT_MAIN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(5.4), card_w - Inches(0.6), Inches(0.5),
                     desc, font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


def build_architecture_slide(prs):
    """Slide 3: System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "System Architecture", "Full-stack overview of the LifeOS platform")

    blocks = [
        ("🖥️  Frontend", SECONDARY, [
            "React + Vite",
            "Inter Typography",
            "Blue/Cyan Gradient Design",
            "Glassmorphism UI",
            "Dark / Light Mode",
            "3-Column Dashboard Layout",
        ]),
        ("⚙️  Backend", PRIMARY, [
            "FastAPI 0.115",
            "Async SQLAlchemy 2.0",
            "17 API Routers",
            "JWT Auth (PyJWT + bcrypt)",
            "Pydantic v2 Validation",
            "Alembic Migrations",
        ]),
        ("🧠  AI Layer", PURPLE, [
            "Groq SDK (Llama 3.3 70B)",
            "AI Chat Assistant",
            "Symptom Analyzer",
            "Nutrition Planner",
            "Fitness Coach",
            "Mental Health Module",
        ]),
    ]

    card_w = Inches(3.6)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, color, items) in enumerate(blocks):
        x = x_start + i * (card_w + gap)
        add_shape(slide, x, Inches(2.0), card_w, Inches(4.8), WHITE,
                  border_color=color, border_width=2)
        add_accent_bar(slide, x, Inches(2.0), card_w, Pt(5), color)
        add_text_box(slide, x + Inches(0.3), Inches(2.2), card_w - Inches(0.6), Inches(0.5),
                     title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.35), Inches(2.9), card_w - Inches(0.7), Inches(3.5),
                        items, font_size=15, color=TEXT_BODY, icon="›")

    # Bottom: Database block
    add_shape(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.45), WHITE,
              border_color=SUCCESS, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(6.91), Inches(11), Inches(0.4),
                 "🗄️  PostgreSQL 16  ·  Async Driver (asyncpg)  ·  Alembic Migrations  ·  Docker + docker-compose",
                 font_size=15, color=SUCCESS, alignment=PP_ALIGN.CENTER)


def build_tech_stack_slide(prs):
    """Slide 4: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "Technology Stack", "Modern tools for a production-grade platform")

    categories = [
        ("Backend Framework", "FastAPI 0.115", "High-performance async Python web framework with automatic OpenAPI docs", PRIMARY),
        ("Database", "PostgreSQL 16 + SQLAlchemy 2.0", "Enterprise-grade relational DB with async ORM (asyncpg driver)", SUCCESS),
        ("Authentication", "JWT + bcrypt + TOTP", "PyJWT tokens, bcrypt password hashing, pyotp for 2FA, Google OAuth", WARNING),
        ("AI / LLM", "Groq SDK — Llama 3.3 70B", "Cloud-hosted large language model for chat, symptoms, nutrition, fitness, mental health", PURPLE),
        ("Frontend", "React + Vite", "Modern SPA with component-based architecture and hot module replacement", SECONDARY),
        ("DevOps", "Docker + docker-compose", "Containerized deployment with PostgreSQL, FastAPI, and Flask Chatbot services", DANGER),
        ("Validation", "Pydantic v2", "Type-safe request/response schemas with automatic validation", PRIMARY_LIGHT),
        ("Migrations", "Alembic", "Version-controlled database schema migrations", RGBColor(0x7C, 0x3A, 0xED)),
    ]

    cols = 2
    card_w = Inches(5.6)
    card_h = Inches(1.1)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    x_start = Inches(0.8)
    y_start = Inches(2.0)

    for i, (cat, tech, desc, color) in enumerate(categories):
        col = i % cols
        row = i // cols
        x = x_start + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)
        add_shape(slide, x, y, card_w, card_h, WHITE, border_color=color, border_width=1)
        add_accent_bar(slide, x, y, Pt(5), card_h, color)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.08), card_w - Inches(0.5), Inches(0.35),
                     f"{cat}:  {tech}", font_size=14, color=color, bold=True)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.5), card_w - Inches(0.5), Inches(0.55),
                     desc, font_size=12, color=TEXT_MUTED)


def build_auth_slide(prs):
    """Slide 5: Authentication & Security."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🔐 Authentication & Security",
                             "Multi-layered identity management")

    cards = [
        ("Email / Password Auth", PRIMARY, [
            "JWT-based session management",
            "bcrypt password hashing",
            "Access + Refresh token flow",
            "Login history tracking",
        ]),
        ("Google OAuth", SECONDARY, [
            "One-click Google Sign-In",
            "Auto-register new users",
            "Seamless 2FA integration",
            "Verified email bypass",
        ]),
        ("Face Verification", PURPLE, [
            "128-float face descriptor matching",
            "Euclidean distance comparison",
            "Client-side face-api.js embedding",
            "Biometric login support",
        ]),
        ("Two-Factor Auth (2FA)", WARNING, [
            "TOTP via authenticator apps",
            "Secure setup with QR URI",
            "Temporary token flow (5 min)",
            "pyotp integration",
        ]),
        ("Email Verification", SUCCESS, [
            "6-digit verification codes",
            "15-minute token expiration",
            "Async email sending (smtplib)",
            "Anti-enumeration protection",
        ]),
        ("Role-Based Access (RBAC)", DANGER, [
            "Patient / Doctor / Admin roles",
            "Role stored on User model",
            "Extensible permission system",
            "Enum-based role management",
        ]),
    ]

    cols = 3
    card_w = Inches(3.7)
    card_h = Inches(2.2)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    x_start = Inches(0.8)
    y_start = Inches(2.0)

    for i, (title, color, items) in enumerate(cards):
        col = i % cols
        row = i // cols
        x = x_start + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, title, items,
                 title_color=color, border_color=color)


def build_healthcare_modules_slide(prs):
    """Slide 6: Core Healthcare Modules."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🩺 Core Healthcare Modules",
                             "Comprehensive patient health management")

    modules = [
        ("Dashboard", SUCCESS, [
            "Health score calculation",
            "Unified summary view",
            "Quick-access widgets",
        ]),
        ("Medical Records", PRIMARY, [
            "Full CRUD operations",
            "File uploads (PDF/images)",
            "AI-powered PDF parsing",
            "AI record comparison",
        ]),
        ("Medicines", WARNING, [
            "Medication tracking",
            "Drug interaction checker",
            "Refill reminders",
        ]),
        ("Appointments", SECONDARY, [
            "Scheduling & management",
            "AI-powered suggestions",
            "Status tracking",
        ]),
        ("Emergency", DANGER, [
            "Emergency contacts & SOS",
            "Medical QR code data",
            "Organ donor status",
        ]),
        ("Family", PURPLE, [
            "Family member profiles",
            "Vaccination records",
            "Medical history sharing",
        ]),
    ]

    cols = 3
    card_w = Inches(3.7)
    card_h = Inches(2.2)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    x_start = Inches(0.8)
    y_start = Inches(2.0)

    for i, (title, color, items) in enumerate(modules):
        col = i % cols
        row = i // cols
        x = x_start + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, title, items,
                 title_color=color, border_color=color)


def build_tracking_slide(prs):
    """Slide 7: Health Tracking & Analytics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "📊 Health Tracking & Analytics",
                             "Data-driven health insights and gamification")

    trackers = [
        ("Health Trackers", PRIMARY, [
            "Water intake (daily glasses)",
            "Sleep quality & duration",
            "Vitals: heart rate, BP, blood sugar",
            "Steps & calories tracking",
            "BMI / BMR calculators",
            "Voice-command data entry",
        ]),
        ("Analytics Engine", SECONDARY, [
            "Health timeline visualization",
            "Interactive charts & graphs",
            "Risk assessment scoring",
            "Predictive health analytics",
            "Trend analysis over time",
            "Data export capabilities",
        ]),
        ("Gamification", WARNING, [
            "Daily health challenges",
            "Progress tracking bars",
            "Streak tracking system",
            "Achievement badges",
            "Task completion rewards",
            "Engagement-driven design",
        ]),
    ]

    card_w = Inches(3.7)
    card_h = Inches(3.6)
    gap = Inches(0.3)
    x_start = Inches(0.8)

    for i, (title, color, items) in enumerate(trackers):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(2.0), card_w, card_h, title, items,
                 title_color=color, border_color=color)

    # Bottom bar — Wearable
    add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.1), WHITE,
              border_color=SUCCESS, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(5.9), Inches(2), Inches(0.4),
                 "⌚ Wearable Integration", font_size=16, color=SUCCESS, bold=True)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(10.8), Inches(0.5),
                 "Fitbit OAuth integration  ·  Auto-sync steps, heart rate & calories  ·  "
                 "Token refresh handling  ·  Connected device management",
                 font_size=14, color=TEXT_BODY)


def build_ai_slide(prs):
    """Slide 8: AI Integration Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🤖 AI Integration — Groq Llama 3.3 70B",
                             "Five specialized AI modules powered by large language models")

    ai_modules = [
        ("💬 AI Chat Assistant", [
            "Context-aware health conversations",
            "User profile injection (age, conditions, allergies)",
            "Chat history persistence",
            "Daily health tips generation",
            "Separate RAG chatbot (Pinecone + Flask)",
        ]),
        ("🔍 Symptom Analyzer", [
            "Multi-symptom analysis",
            "Duration & severity weighting",
            "Condition probability scoring",
            "Specialist recommendations",
            "Hardcoded fallback dictionary",
        ]),
        ("🥗 AI Nutritionist", [
            "Personalized meal planning",
            "Macro tracking & stats",
            "Nutritional recommendations",
            "Diet optimization",
            "Calorie-aware suggestions",
        ]),
        ("💪 AI Fitness Coach", [
            "Custom workout generation",
            "Weekly fitness planning",
            "Step tracking integration",
            "Fitness statistics",
            "Voice intent logging",
        ]),
        ("🧘 AI Mental Health", [
            "Mood tracking & logging",
            "AI-assisted journaling",
            "Stress level analysis",
            "Psychological screening",
            "Mental wellness tips",
        ]),
    ]

    card_w = Inches(2.15)
    card_h = Inches(3.6)
    gap = Inches(0.2)
    x_start = Inches(0.8)

    colors = [PRIMARY, DANGER, SUCCESS, WARNING, PURPLE]

    for i, ((title, items), color) in enumerate(zip(ai_modules, colors)):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(2.0), card_w, card_h, title, items,
                 title_color=color, border_color=color)

    # Bottom: PDF AI
    add_shape(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0), WHITE,
              border_color=PRIMARY, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(5.98), Inches(10.8), Inches(0.35),
                 "📄 AI PDF Parsing — Medical Report Intelligence",
                 font_size=16, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(1.1), Inches(6.35), Inches(10.8), Inches(0.5),
                 "PyMuPDF text extraction  →  Groq AI metric extraction  →  Auto health-entry creation  →  "
                 "AI-powered record comparison between reports",
                 font_size=14, color=TEXT_BODY)


def build_api_modules_slide(prs):
    """Slide 9: API Endpoints Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "📡 API Modules — 80+ Endpoints",
                             "RESTful API architecture with Swagger documentation")

    modules = [
        ("Authentication",   "/api/v1/auth",         "register, login, refresh, logout, 2FA, OAuth, face"),
        ("Users",            "/api/v1/users",         "profile, settings, data export"),
        ("Dashboard",        "/api/v1/dashboard",     "summary, health-score"),
        ("Medical Records",  "/api/v1/records",       "CRUD + file upload + AI summary + compare"),
        ("Medicines",        "/api/v1/medicines",     "CRUD + interactions + refill reminders"),
        ("Appointments",     "/api/v1/appointments",  "CRUD + AI prep notes"),
        ("Emergency",        "/api/v1/emergency",     "contacts, SOS, QR data, organ donor"),
        ("Family",           "/api/v1/family",        "members + vaccination records"),
        ("Trackers",         "/api/v1/trackers",      "water, sleep, health metrics, BMI, voice"),
        ("Expenses",         "/api/v1/expenses",      "CRUD + category summaries"),
        ("Challenges",       "/api/v1/challenges",    "progress, complete, streak, badges"),
        ("Analytics",        "/api/v1/analytics",     "timeline, graphs, risk, predictions"),
        ("AI Chat",          "/api/v1/ai/chat",       "chat, history, daily tips"),
        ("AI Symptoms",      "/api/v1/ai/symptoms",   "multi-symptom analysis"),
        ("AI Nutrition",     "/api/v1/ai/nutrition",  "meal plan, macro stats"),
        ("AI Fitness",       "/api/v1/ai/fitness",    "workouts, weekly plan, step stats"),
    ]

    # Table header
    add_shape(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.45), PRIMARY)
    add_text_box(slide, Inches(1.0), Inches(1.92), Inches(3), Inches(0.4),
                 "Module", font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.2), Inches(1.92), Inches(3.2), Inches(0.4),
                 "Prefix", font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.5), Inches(1.92), Inches(4.5), Inches(0.4),
                 "Key Endpoints", font_size=14, color=WHITE, bold=True)

    # Table rows
    row_h = Inches(0.3)
    y_start = Inches(2.4)
    for i, (module, prefix, endpoints) in enumerate(modules):
        y = y_start + i * (row_h + Pt(2))
        bg = WHITE if i % 2 == 0 else BG_SECONDARY
        add_shape(slide, Inches(0.8), y, Inches(11.5), row_h, bg)
        add_text_box(slide, Inches(1.0), y, Inches(3), row_h,
                     module, font_size=12, color=TEXT_MAIN, bold=True)
        add_text_box(slide, Inches(4.2), y, Inches(3.2), row_h,
                     prefix, font_size=11, color=PRIMARY)
        add_text_box(slide, Inches(7.5), y, Inches(4.5), row_h,
                     endpoints, font_size=11, color=TEXT_MUTED)


def build_database_slide(prs):
    """Slide 10: Database Schema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🗄️ Database Schema",
                             "11 SQLAlchemy ORM models with PostgreSQL 16")

    tables = [
        ("users", "Core user identity", "id, email, hashed_password, role, is_active, is_verified, 2FA, face_descriptor"),
        ("user_profiles", "Extended profile data", "name, age, blood_type, conditions, allergies, fitbit_tokens, avatar"),
        ("medical_records", "Patient medical records", "title, category, doctor, hospital, date, findings, file_path"),
        ("health_entries", "Tracked health metrics", "category (BP, HR, sugar, weight, steps, calories), value, recorded_at"),
        ("water_intake", "Daily hydration log", "user_id, date, glasses — unique constraint per (user, date)"),
        ("sleep_entries", "Sleep tracking", "date, hours, quality, bedtime, wake_time"),
        ("chat_messages", "AI conversation history", "user_id, role (user/assistant), content, module, created_at"),
        ("medical_expenses", "Cost tracking", "description, category (medicine, doctor, tests, insurance), amount, date"),
        ("family_members", "Family profiles", "name, relation, age, blood_type, conditions, medications"),
        ("appointments", "Doctor appointments", "doctor, specialty, hospital, date, time, status, ai_prep_notes"),
        ("login_history", "Login audit trail", "user_id, ip_address, user_agent, created_at"),
    ]

    row_h = Inches(0.42)
    y_start = Inches(2.0)
    # Header
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.4), PRIMARY)
    add_text_box(slide, Inches(1.0), Inches(1.82), Inches(2.5), Inches(0.35),
                 "Table", font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.5), Inches(1.82), Inches(2.5), Inches(0.35),
                 "Purpose", font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(6.0), Inches(1.82), Inches(6), Inches(0.35),
                 "Key Columns", font_size=13, color=WHITE, bold=True)

    for i, (table, purpose, cols) in enumerate(tables):
        y = y_start + Inches(0.2) + i * row_h
        bg = WHITE if i % 2 == 0 else BG_SECONDARY
        add_shape(slide, Inches(0.8), y, Inches(11.5), row_h - Pt(3), bg)
        add_text_box(slide, Inches(1.0), y + Pt(1), Inches(2.3), row_h,
                     table, font_size=12, color=PRIMARY, bold=True)
        add_text_box(slide, Inches(3.5), y + Pt(1), Inches(2.3), row_h,
                     purpose, font_size=11, color=TEXT_BODY)
        add_text_box(slide, Inches(6.0), y + Pt(1), Inches(6), row_h,
                     cols, font_size=10, color=TEXT_MUTED)


def build_design_slide(prs):
    """Slide 11: UI/UX Design System."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🎨 Design System",
                             "Modern UI 2.1 aesthetic — Glassmorphism + Gradients")

    # Color palette section
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(3), Inches(0.4),
                 "Color Palette", font_size=18, color=TEXT_MAIN, bold=True)

    palette = [
        ("#2563EB", "Primary",   RGBColor(0x25, 0x63, 0xEB)),
        ("#06B6D4", "Secondary", RGBColor(0x06, 0xB6, 0xD4)),
        ("#22C55E", "Success",   RGBColor(0x22, 0xC5, 0x5E)),
        ("#F59E0B", "Warning",   RGBColor(0xF5, 0x9E, 0x0B)),
        ("#EF4444", "Danger",    RGBColor(0xEF, 0x44, 0x44)),
        ("#F0F4FF", "Light BG",  RGBColor(0xD0, 0xD8, 0xEE)),  # slightly darker for visibility
    ]

    x = Inches(0.8)
    for hex_val, name, color in palette:
        swatch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(1.7), Inches(1.0))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.color.rgb = CARD_BORDER
        swatch.line.width = Pt(1)
        swatch.adjustments[0] = 0.15
        add_text_box(slide, x, Inches(3.55), Inches(1.7), Inches(0.3),
                     name, font_size=12, color=TEXT_MAIN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.8), Inches(1.7), Inches(0.25),
                     hex_val, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
        x += Inches(1.85)

    # Design principles cards
    principles = [
        ("Typography", PRIMARY, [
            "Primary: Inter font family",
            "Headings: Bold 700–800",
            "Body: Regular 400–500",
            "Base: 16px (1rem)",
        ]),
        ("Glassmorphism", SECONDARY, [
            "blur(16px) backdrop filter",
            "Semi-transparent panels",
            "Soft layered shadows",
            "rgba overlays for depth",
        ]),
        ("Layout", SUCCESS, [
            "3-column CSS Grid",
            "Sidebar: 260px",
            "Main content: 1fr",
            "Right panel: 320px",
            "Max width: 1440px",
        ]),
        ("Theming", WARNING, [
            "CSS custom properties",
            "Light & Dark modes",
            "Isolated theme contexts",
            "data-theme attribute toggle",
        ]),
    ]

    card_w = Inches(2.7)
    gap = Inches(0.27)
    x_start = Inches(0.8)
    for i, (title, color, items) in enumerate(principles):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(4.3), card_w, Inches(2.7), title, items,
                 title_color=color, border_color=color)


def build_deployment_slide(prs):
    """Slide 12: Deployment & DevOps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🐳 Deployment & Infrastructure",
                             "Containerized deployment with Docker")

    # Docker compose services
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4),
                 "Docker Compose Services", font_size=20, color=TEXT_MAIN, bold=True)

    services = [
        ("lifeos_db", "PostgreSQL 16", "Primary database  ·  Port 5432  ·  Persistent volumes", SECONDARY),
        ("lifeos_app", "FastAPI Application", "Main backend  ·  Port 8000  ·  Auto Alembic migrations", PRIMARY),
        ("lifeos_chatbot", "Flask RAG Chatbot", "Pinecone vector DB  ·  Separate AI chat service", PURPLE),
    ]

    for i, (name, tech, desc, color) in enumerate(services):
        y = Inches(2.6) + i * Inches(1.15)
        add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.95), WHITE,
                  border_color=color, border_width=1)
        add_accent_bar(slide, Inches(0.8), y, Pt(5), Inches(0.95), color)
        add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(5), Inches(0.35),
                     f"{name}  —  {tech}", font_size=15, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.48), Inches(5), Inches(0.4),
                     desc, font_size=13, color=TEXT_MUTED)

    # Environment variables
    add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5), Inches(0.4),
                 "Environment Variables", font_size=20, color=TEXT_MAIN, bold=True)

    env_vars = [
        ("DATABASE_URL", "PostgreSQL connection string"),
        ("SECRET_KEY", "JWT signing key"),
        ("GROQ_API_KEY", "AI model API key"),
        ("GOOGLE_CLIENT_ID", "OAuth client ID"),
        ("FITBIT_CLIENT_ID", "Wearable integration"),
        ("CORS_ORIGINS", "Allowed frontend origins"),
        ("DEBUG", "Development mode flag"),
        ("MAX_FILE_SIZE_MB", "Upload size limit"),
    ]

    y_start = Inches(2.6)
    for i, (var, desc) in enumerate(env_vars):
        y = y_start + i * Inches(0.5)
        bg = WHITE if i % 2 == 0 else BG_SECONDARY
        add_shape(slide, Inches(7.0), y, Inches(5.5), Inches(0.42), bg,
                  border_color=CARD_BORDER, border_width=1)
        add_text_box(slide, Inches(7.2), y + Pt(1), Inches(2.5), Inches(0.35),
                     var, font_size=12, color=PRIMARY, bold=True)
        add_text_box(slide, Inches(9.7), y + Pt(1), Inches(2.5), Inches(0.35),
                     desc, font_size=11, color=TEXT_MUTED)


def build_future_slide(prs):
    """Slide 13: Future Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)
    add_section_header_shape(slide, "🚀 Future Roadmap",
                             "Planned enhancements and next-phase features")

    roadmap = [
        ("Phase 1 — Security Hardening", PRIMARY, [
            "Server-side token revocation / blacklist",
            "Robust RBAC middleware with @requires_role",
            "Backup recovery codes for 2FA",
            "Anti-spoofing / liveness for face auth",
        ]),
        ("Phase 2 — Real-Time Features", SECONDARY, [
            "WebSocket-based live notifications",
            "Push notifications (FCM / APNS)",
            "In-app notification system",
            "Real-time health alerts",
        ]),
        ("Phase 3 — Platform Expansion", SUCCESS, [
            "FHIR / HL7 health data standards",
            "Cloud storage (S3 / Cloudinary)",
            "Payment integration for billing",
            "Calendar sync (Google Calendar)",
        ]),
        ("Phase 4 — AI Evolution", PURPLE, [
            "Full RAG pipeline in main app",
            "Multi-turn conversation context",
            "OCR for image-based PDFs",
            "Predictive health modeling",
        ]),
    ]

    card_w = Inches(2.7)
    card_h = Inches(3.2)
    gap = Inches(0.27)
    x_start = Inches(0.8)

    for i, (title, color, items) in enumerate(roadmap):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(2.0), card_w, card_h, title, items,
                 title_color=color, border_color=color)


def build_thank_you_slide(prs):
    """Slide 14: Thank You / Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)

    # Decorative circles — light pastels
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CIRCLE_LIGHT
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = CIRCLE_CYAN
    circle2.line.fill.background()

    add_text_box(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(1.0),
                 "Thank You", font_size=56, color=TEXT_MAIN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(5.5), Inches(3.1), Inches(2.3), Pt(5), SECONDARY)

    add_text_box(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.7),
                 "LifeOS — AI-Powered Healthcare Operating System",
                 font_size=22, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(4.3), Inches(13.333), Inches(0.5),
                 "Built with FastAPI  ·  React  ·  PostgreSQL  ·  Groq Llama 3.3  ·  Docker",
                 font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Links card
    add_shape(slide, Inches(4.5), Inches(5.2), Inches(4.3), Inches(1.4), WHITE,
              border_color=PRIMARY, border_width=1)
    add_text_box(slide, Inches(4.8), Inches(5.3), Inches(3.7), Inches(0.35),
                 "📚  Quick Links", font_size=16, color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), Inches(5.7), Inches(3.7), Inches(0.3),
                 "Swagger:  localhost:8000/docs", font_size=13, color=TEXT_BODY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), Inches(6.0), Inches(3.7), Inches(0.3),
                 "Frontend:  localhost:5173", font_size=13, color=TEXT_BODY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), Inches(6.3), Inches(3.7), Inches(0.3),
                 "ReDoc:  localhost:8000/redoc", font_size=13, color=TEXT_BODY,
                 alignment=PP_ALIGN.CENTER)


# ──────────────────────────── MAIN ────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    build_title_slide(prs)           # 1
    build_overview_slide(prs)        # 2
    build_architecture_slide(prs)    # 3
    build_tech_stack_slide(prs)      # 4
    build_auth_slide(prs)            # 5
    build_healthcare_modules_slide(prs)  # 6
    build_tracking_slide(prs)        # 7
    build_ai_slide(prs)              # 8
    build_api_modules_slide(prs)     # 9
    build_database_slide(prs)        # 10
    build_design_slide(prs)          # 11
    build_deployment_slide(prs)      # 12
    build_future_slide(prs)          # 13
    build_thank_you_slide(prs)       # 14

    # Add page numbers (skip title and thank you)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        add_page_number(slide, i + 1, total)

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "LifeOS_Healthcare_AI_Presentation.pptx")
    prs.save(output_path)
    print(f"\n✅ Presentation saved successfully!")
    print(f"📁 Path: {output_path}")
    print(f"📊 Total slides: {total}")
    print(f"\nSlide Overview:")
    slide_names = [
        "Title Slide", "Project Overview", "System Architecture",
        "Technology Stack", "Authentication & Security",
        "Core Healthcare Modules", "Health Tracking & Analytics",
        "AI Integration", "API Modules (80+ endpoints)",
        "Database Schema", "Design System", "Deployment & DevOps",
        "Future Roadmap", "Thank You"
    ]
    for i, name in enumerate(slide_names):
        print(f"  {i+1:2d}. {name}")


if __name__ == "__main__":
    main()
